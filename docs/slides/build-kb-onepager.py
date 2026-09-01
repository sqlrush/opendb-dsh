#!/usr/bin/env python3
"""opendb-harness 客户知识库一页 PPT（user 2026-09-01 要的）：
   把客户既有的几十篇规范 / 上千条工单 / 上百份故障总结，自动构建成"真懂客户"的知识库；
   并说清三类存储引擎（关系型 / 向量 / 图）各自的角色与相互配合。
   刻意只写存储**类型**，不出现任何产品名。

   /usr/bin/python3 docs/slides/build-kb-onepager.py  →  docs/slides/opendb-harness-知识库一页.pptx
   预览：qlmanage -t -s 2000 -o /tmp/ql <pptx>（macOS QuickLook 渲染成 PNG）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

OUT = Path(__file__).with_name('opendb-harness-知识库一页.pptx')
FONT = 'PingFang SC'
INK, SUB, DIM, BLUE, LINE = '0F1115', '61666B', '81858C', '4176E6', 'D9DCE1'
GREEN, AMBER, PURPLE, TEAL, RED = '3FA552', 'C9862D', '8B6BE0', '2FA79A', 'D6604D'
FILL = {'src': 'FFFFFF', 'ingest': 'EEF3FF', 'extract': 'E8F5EC', 'rel': 'F2F3F5',
        'vec': 'FDF0E3', 'graph': 'F3EEFC', 'search': 'EEF3FF', 'use': 'F7F8FA'}
EDGE = {'src': 'C9CED6', 'ingest': BLUE, 'extract': GREEN, 'rel': '8A9099',
        'vec': AMBER, 'graph': PURPLE, 'search': BLUE, 'use': 'B8BCC4'}

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def rgb(h):
    return RGBColor.from_string(h)


def text(x, y, w, h, runs, size=10, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        parts = [(para, {})] if isinstance(para, str) else para
        for t, st in parts:
            r = p.add_run()
            r.text = t
            r.font.name = FONT
            r.font.size = Pt(st.get('size', size))
            r.font.bold = st.get('bold', bold)
            r.font.color.rgb = rgb(st.get('color', color))
    return tb


def box(x, y, w, h, kind, title, badge=None, lines=(), title_size=11.5, body_size=9, radius=0.08, badge_size=9):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(FILL[kind])
    shp.line.color.rgb = rgb(EDGE[kind]); shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    shp.text_frame.text = ''
    text(x + 0.08, y + 0.04, w - 0.16, 0.3,
         [[(title, {'bold': True, 'size': title_size, 'color': INK})]
          + ([('  ' + badge, {'bold': True, 'size': badge_size, 'color': EDGE[kind]})] if badge else [])])
    if lines:
        text(x + 0.08, y + 0.32, w - 0.16, h - 0.36, ['· ' + ln for ln in lines], size=body_size, color=SUB, line_spacing=1.12)
    return shp


def arrow(x1, y1, x2, y2, color=SUB, dashed=False, both=False, width=1.25, plain=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = rgb(color); c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if plain:
        return c
    if both:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def label(x, y, w, t, color=DIM, size=8, align=PP_ALIGN.LEFT):
    tb = text(x, y, w, 0.16, [t], size=size, color=color, align=align, anchor=MSO_ANCHOR.MIDDLE)
    tb.text_frame.margin_top = tb.text_frame.margin_bottom = 0
    return tb


# ── 标题 ──────────────────────────────────────────────────────────────
text(0.5, 0.26, 12.3, 0.5, [[('客户知识库：从"存着"到"真懂"', {'bold': True, 'size': 24, 'color': INK}),
                             ('   把客户既有文本自动构建成可被平台使用的知识资产', {'size': 14, 'color': SUB})]])
text(0.5, 0.76, 12.3, 0.3, ['几十篇使用规范 + 上千条问题处理工单 + 上百份故障分析总结 → 结构化 → 三类存储各司其职 → 每一条诊断结论旁都能说出"贵司规范怎么说、上次怎么处理"'],
     size=11, color=BLUE)

L, W = 0.5, 8.15

# ── ① 客户既有资料 ───────────────────────────────────────────────────
sw = (W - 0.24) / 3
for i, (t, b, n) in enumerate([
    ('使用规范', '运维/变更/安全规范、准入标准', '几十篇'),
    ('问题处理工单', '现象、处置动作、验证与耗时', '上千条'),
    ('故障分析总结', '根因、影响面、复发防范', '上百份'),
]):
    x = L + i * (sw + 0.12)
    box(x, 1.18, sw, 0.62, 'src', t, badge=n, lines=[b], title_size=11, body_size=8.4, badge_size=8.5)
text(L, 1.02, W, 0.16, ['① 客户既有资料（多来源、多格式、无统一结构）'], size=8.5, color=DIM)

# ── ② 摄入与治理 ─────────────────────────────────────────────────────
box(L, 2.02, W, 0.60, 'ingest', '② 摄入与治理', badge='自动化管线 · 增量',
    lines=['格式归一（文档 / 表格 / 工单导出 / 邮件）→ 语义切块；每篇强制打元数据：类型 · 适用引擎与环境 · 生效时间 · 版本 · 密级',
           '同源资料重灌产出新版本而非覆盖——报告引用永远能追溯到"当时依据的是哪一版规范"'],
    title_size=11.5, body_size=8.6)

# ── ③ 结构化抽取 ─────────────────────────────────────────────────────
box(L, 2.84, W, 1.00, 'extract', '③ 结构化抽取：把文本变成"可判定/可复用"的知识', badge='模型抽取 + 人工确认后才可被引用',
    title_size=11.5, body_size=8.6, badge_size=8.5)
ew = (W - 0.36) / 3
for i, (t, b, c) in enumerate([
    ('规范 → 条款卡', '约束对象 · 条件 · 要求值 · 强制/建议', GREEN),
    ('工单 → 处置卡', '现象 · 动作 · 验证方式 · 耗时 · 涉及对象', GREEN),
    ('故障总结 → 案例卡', '现象 · 影响 · 根因 · 处置 · 防复发', GREEN),
]):
    x = L + 0.12 + i * (ew + 0.06)
    text(x, 3.24, ew, 0.5, [[(t, {'bold': True, 'size': 9.5, 'color': INK})], b], size=8.2, color=SUB, line_spacing=1.12)

# ── ④ 三类存储各司其职 ───────────────────────────────────────────────
text(L, 3.94, W, 0.16, ['④ 三类存储各司其职：真相只有一份，能力互补，任一层不可用都能降级'], size=8.5, color=DIM)
by, bh = 4.12, 1.36
GAP = 0.30                      # 缝要够宽，双向箭头才看得见（0.16 的缝里箭头几乎被框线吃掉）
bw = (W - 2 * GAP) / 3
box(L, by, bw, bh, 'rel', '关系型数据库', badge='唯一真相',
    lines=['文档 · 版本 · 条款卡 / 处置卡 / 案例卡',
           '记忆系统：平台在客户环境里做过什么',
           '引用台账：哪条结论引用了哪一版哪一条',
           '精确过滤：引擎 · 环境 · 生效期 · 权限'],
    title_size=11, body_size=8.3, badge_size=8.5)
box(L + bw + GAP, by, bw, bh, 'vec', '向量数据库', badge='语义召回',
    lines=['切块与卡片的向量索引，规模化近邻检索',
           '解决"说法不同、意思相同"的召回',
           '客户说"连接打满" ↔ 平台判"连接超阈值"',
           '不存业务事实——只加速，可随时重建'],
    title_size=11, body_size=8.3, badge_size=8.5)
box(L + 2 * (bw + GAP), by, bw, bh, 'graph', '图数据库', badge='关系推理',
    lines=['实体：对象 · 现象 · 根因 · 处置动作 · 条款',
           '边：现象→根因→处置、条款→约束对象',
           '案例→引用条款、故障→涉及节点',
           '多跳串联：一条结论牵出规范与历史全链'],
    title_size=11, body_size=8.3, badge_size=8.5)
# 三库之间的配合：缝里一根双向箭头 + 缝下的说明
mid = by + bh / 2
arrow(L + bw, mid, L + bw + GAP, mid, both=True, color=AMBER, width=1.3)
arrow(L + 2 * bw + GAP, mid, L + 2 * bw + 2 * GAP, mid, both=True, color=PURPLE, width=1.3)
label(L + bw + GAP / 2 - 0.75, by + bh + 0.03, 1.5, '同步向量 / 回表取原文', color=AMBER, size=7.5, align=PP_ALIGN.CENTER)
label(L + 2 * bw + 1.5 * GAP - 0.8, by + bh + 0.03, 1.6, '抽实体与边 / 回表取证据', color=PURPLE, size=7.5, align=PP_ALIGN.CENTER)

# ── ⑤ 混合检索编排 ───────────────────────────────────────────────────
box(L, 5.72, W, 0.58, 'search', '⑤ 混合检索编排（由平台按"发现"发起，不靠模型自己想起来查）', badge='语义 + 词法 + 结构过滤 + 图扩展 → 重排',
    lines=['检索键 = 规则码 + 对象 + 现象词：语义召回近义说法、词法精确命中错误码与对象名、按引擎/环境/生效期过滤、再沿图扩展到根因与历史处置'],
    title_size=11, body_size=8.6, badge_size=8.5)

# ── ⑥ 落到平台的确定性结论上 ─────────────────────────────────────────
box(L, 6.42, W, 0.86, 'use', '⑥ 用在哪：贴到平台每一条确定性结论旁', badge='引用必须有出处 · 查不到就写"无对应规范"，绝不编造',
    title_size=11, badge_size=8.5)
uw = (W - 0.36) / 4
for i, (t, b) in enumerate([
    ('规范对照', '结论旁引用条款原文（参考不改判）'),
    ('建议本地化', '按客户流程改写处置建议'),
    ('阈值对齐', '口径差异清单 → 人确认后生效'),
    ('案例复用', '相似历史故障的处置与耗时'),
]):
    x = L + 0.12 + i * (uw + 0.04)
    text(x, 6.76, uw, 0.46, [[(t, {'bold': True, 'size': 9, 'color': INK})], b], size=8, color=SUB, line_spacing=1.1)

# 主流程箭头
cx = L + W / 2
arrow(cx, 1.80, cx, 2.02, color=BLUE)
arrow(cx, 2.62, cx, 2.84, color=GREEN)
arrow(cx, 3.84, cx, 4.12, color=SUB)
arrow(cx, by + bh, cx, 5.72, color=BLUE)
arrow(cx, 6.30, cx, 6.42, color=SUB)

# ── 右栏：怎么做到"真懂客户" ─────────────────────────────────────────
R, RW = 8.95, 3.9
panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(R), Inches(1.18), Inches(RW), Inches(6.10))
panel.adjustments[0] = 0.04
panel.fill.solid(); panel.fill.fore_color.rgb = rgb('FFFFFF')
panel.line.color.rgb = rgb(LINE); panel.line.width = Pt(1)
panel.shadow.inherit = False
text(R + 0.15, 1.26, RW - 0.3, 0.35, [[('为什么这样能"真懂客户"', {'bold': True, 'size': 14, 'color': INK})]])
points = [
    ('不止"能搜"，而是"能对账"',
     '文本变成条款卡后，可与平台内置的确定性规则逐条比对：客户要求平台已在判但阈值不同 → 给出差异清单；客户有要求平台没规则 → 能力缺口；平台在判客户没写 → 建议补进规范。这一步把知识库从"资料检索"变成"按客户口径工作"。'),
    ('两套经验双轮驱动',
     '知识 = 客户自己沉淀的规范与案例；记忆 = 平台在这套环境里做过什么、上次结论是什么。两者同库同源、互为佐证——既懂客户的规矩，也记得这套库的脾气。'),
    ('三类存储互补且互为降级',
     '关系型存真相与权限、向量做语义召回、图做多跳串联；三者共用同一套标识，向量或图不可用时自动退回关系型，检索质量下降但结论不出错。'),
    ('越用越懂：缺口驱动地长大',
     '报告里"无对应规范"的发现自动排队，提示补充哪类资料；DBA 对每条引用点"有用/无关"回写排序权重。冷启动只需先灌最核心的几十篇，而不是一次性倒进所有历史文档。'),
]
y = 1.70
for i, (t, body) in enumerate(points):
    num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(R + 0.18), Inches(y + 0.02), Inches(0.3), Inches(0.3))
    num.fill.solid(); num.fill.fore_color.rgb = rgb(BLUE); num.line.fill.background(); num.shadow.inherit = False
    tf = num.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i + 1); r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = rgb('FFFFFF'); r.font.name = FONT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(R + 0.58, y - 0.02, RW - 0.75, 0.34, [[(t, {'bold': True, 'size': 11.5, 'color': INK})]])
    text(R + 0.58, y + 0.30, RW - 0.75, 1.15, [body], size=9.2, color=SUB, line_spacing=1.18)
    y += 1.44

text(0.5, 7.32, 8.2, 0.18, ['图例：实线箭头 = 主流程；框间双向箭头 = 存储层之间的同步与回表；三类存储只写类型，不绑定具体产品'], size=7.5, color=DIM)
text(9.0, 7.32, 3.9, 0.18, ['opendb-harness · 客户知识库一页 · 2026-09'], size=7.5, color=DIM, align=PP_ALIGN.RIGHT)

prs.save(OUT)
print('saved', OUT)
