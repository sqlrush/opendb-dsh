#!/usr/bin/env python3
"""opendb-harness 架构一页 PPT（user 2026-08-28 临时要的）：python-pptx 画 k8s 多角色 Pod 架构图 + 架构特点。
   python3 docs/slides/build-arch-onepager.py  →  docs/slides/opendb-harness-架构一页.pptx
   预览：qlmanage -t -s 2000 -o /tmp/ql <pptx>（macOS QuickLook 渲染成 PNG）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

OUT = Path(__file__).with_name('opendb-harness-架构一页.pptx')
FONT = 'PingFang SC'
INK, SUB, DIM, BLUE, LINE = '0F1115', '61666B', '81858C', '4176E6', 'D9DCE1'
FILL = {'user': 'FFFFFF', 'host': 'EEF3FF', 'runtime': 'E8F5EC', 'collector': 'FAF3E5', 'state': 'F2F3F5', 'db': 'FDF0E3', 'model': 'F7F8FA'}
EDGE = {'user': 'C9CED6', 'host': '4176E6', 'runtime': '3FA552', 'collector': 'C9862D', 'state': 'B8BCC4', 'db': 'E07A1F', 'model': 'B8BCC4'}

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def rgb(h):
    return RGBColor.from_string(h)


def text(x, y, w, h, runs, size=10, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """runs: list of paragraphs; each paragraph = str or list of (text, {bold,color,size})."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        parts = [(para, {})] if isinstance(para, str) else para
        for t, st in parts:
            r = p.add_run()
            r.text = t
            r.font.name = FONT
            r.font.size = Pt(st.get('size', size))
            r.font.bold = st.get('bold', bold)
            r.font.color.rgb = rgb(st.get('color', color))
    return tb


def box(x, y, w, h, kind, title, badge=None, lines=(), title_size=11.5, body_size=9, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(FILL[kind])
    shp.line.color.rgb = rgb(EDGE[kind]); shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    shp.text_frame.text = ''
    # 标题 + 副本徽章
    text(x + 0.08, y + 0.04, w - 0.16, 0.3, [[(title, {'bold': True, 'size': title_size, 'color': INK})] + ([('  ' + badge, {'bold': True, 'size': 9, 'color': EDGE[kind]})] if badge else [])])
    if lines:
        text(x + 0.08, y + 0.34, w - 0.16, h - 0.38, ['· ' + ln for ln in lines], size=body_size, color=SUB, line_spacing=1.12)
    return shp


def arrow(x1, y1, x2, y2, color=SUB, dashed=False, both=False, width=1.25, plain=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = rgb(color); c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if plain:
        return c
    if both:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def label(x, y, w, t, color=DIM, size=8, align=PP_ALIGN.LEFT, rotation=0.0):
    """框与框之间 0.2in 的缝里放的小字：无底色、无边框，高度压到 0.16，不压框线"""
    tb = text(x, y, w, 0.16, [t], size=size, color=color, align=align, anchor=MSO_ANCHOR.MIDDLE)
    tb.text_frame.margin_top = tb.text_frame.margin_bottom = 0
    if rotation:
        tb.rotation = rotation
    return tb


# ── 标题 ──────────────────────────────────────────────────────────────
text(0.5, 0.28, 12.3, 0.5, [[('opendb-harness 架构总览', {'bold': True, 'size': 24, 'color': INK}), ('   基于 Kubernetes 的多角色 Pod 架构', {'size': 14, 'color': SUB})]])
text(0.5, 0.78, 12.3, 0.3, ['所有角色皆为可水平扩展的 Pod · 应对大型数据库集群管理 · 支撑数千人同时访问 · 上下文与记忆在 Pod 之间无缝切换'], size=11, color=BLUE)

# ── 架构图（左）──────────────────────────────────────────────────────
L, W = 0.5, 8.1
# 用户层
box(L, 1.2, W, 0.62, 'user', '用户层', badge='数千并发用户', lines=['浏览器（DBA / 开发 / 值班）通过 Ingress · Service 负载均衡接入，HTTP 与 WebSocket 可落到任意 Host 副本'], body_size=9)
# Host
WM = 7.35   # Host / Runtime / Collector / 模型 这四层收窄，右侧留出 0.75in 的状态总线通道
box(L, 2.02, WM, 1.08, 'host', 'Host Pod（接入与调度层）', badge='Deployment ×N · 无状态 · KEDA 按连接数扩缩',
    lines=['会话接入：HTTP / WebSocket、认证、排队投影（提问入队、可编辑 / 插队 / 撤回）',
           '任务调度引擎（cron / 手动 / 事件）、任务大盘与插件 UI 服务、会话日志镜像回放',
           '跨副本扇出：任一副本收到的事件经 PG NOTIFY 同步到所有副本——用户落在哪个副本都一致'])
# Runtime
box(L, 3.3, WM, 1.18, 'runtime', 'Runtime Pod（执行层）', badge='Deployment ×N · KEDA 按队列深度扩缩（2 → 20+）',
    lines=['从 PostgreSQL 队列认领会话轮次，编排大模型推理与工具调用，写会话日志（单写者）',
           '只读诊断工具：SQL 执行 / EXPLAIN / 指标 / 健康 · Top SQL · WDR · DDL 采集器，结果存档供大盘直读',
           '心跳所有权栅栏 + 轮次看门狗 + 失败重投：Pod 滚动或故障时提问不丢、不重复执行'])
# Collector + 模型服务
box(L, 4.68, 3.9, 0.86, 'collector', 'Collector Pod（采集层）', badge='×N · 可按节点分片',
    lines=['每分钟采集指标（TPS / 连接 / 等待 / 缓存…）与数据字典快照', '写入 TimescaleDB：趋势曲线 · 健康判定 · DDL 变更追溯'], body_size=8.6)
box(L + 4.1, 4.68, 3.25, 0.86, 'model', '大模型服务', badge='外部 API / 本机推理',
    lines=['Runtime 统一编排调用，可按会话 / 任务切换模型', '模型只做归因与建议，数字来自确定性采集'], body_size=8.6)
# 状态层
box(L, 5.74, W, 0.9, 'state', '共享状态层（StatefulSet / 持久卷）', badge='所有会话、记忆、队列、指标的唯一真相')
sx = L + 0.1
for i, (t, b) in enumerate([
    ('PostgreSQL + TimescaleDB', '会话日志 · 排队 · 任务 · 指标 · 记忆(pgvector) · 知识 · 阈值 · 采集存档'),
    ('Redis', 'kv 状态 · 会话注册表 · 版本键（跨副本）'),
    ('MinIO', '附件 · 溢写 · 大结果'),
    ('Qdrant', '记忆 / 知识向量检索加速'),
]):
    ws = [2.9, 1.75, 1.45, 1.6][i]
    text(sx, 6.08, ws, 0.6, [[(t, {'bold': True, 'size': 9.5, 'color': INK})], b], size=8.2, color=SUB, line_spacing=1.1)
    sx += ws + 0.1
# 被管数据库集群
box(L, 6.82, W, 0.46, 'db', '被管数据库集群', badge='openGauss / PostgreSQL 节点 ×N（成百上千）',
    lines=[], title_size=11)
text(L + 4.55, 6.84, W - 4.65, 0.42, ['只读账号访问，能做什么由数据库侧权限决定；Runtime 诊断 / Collector 采集分别连接各节点'], size=8.2, color=SUB, anchor=MSO_ANCHOR.MIDDLE)

# 箭头
cx = L + WM / 2
arrow(cx, 1.82, cx, 2.02, both=True, color=BLUE)                       # 用户 ↔ Host
label(cx + 0.12, 1.84, 1.9, 'HTTP / WebSocket')
arrow(cx, 3.10, cx, 3.30, dashed=True, both=True, color=SUB)          # Host ↔ Runtime（经 PG 队列，不直连）
label(cx + 0.12, 3.12, 3.4, '不直连：经 PostgreSQL 队列会合 / 日志镜像')
arrow(L + 1.8, 4.48, L + 1.8, 4.68, color=EDGE['collector'])          # Runtime → Collector（采集器调用 / 结果存档）
label(L + 1.92, 4.50, 2.2, '采集器调用 / 结果存档')
arrow(L + 5.72, 4.48, L + 5.72, 4.68, both=True, color=SUB)           # Runtime ↔ 模型
label(L + 5.84, 4.50, 1.4, '推理调用')
arrow(L + 1.8, 5.54, L + 1.8, 5.74, color=SUB)                         # Collector → 状态层（指标 / 字典）
label(L + 1.92, 5.56, 1.8, '指标 · 字典快照')
arrow(cx, 6.64, cx, 6.82, both=True, color=EDGE['db'])                 # 状态层 ↔ 被管 DB（示意：采集/诊断连接）
label(cx + 0.12, 6.65, 3.0, '只读诊断 / 指标采集连接（经 Runtime / Collector）')
# Host / Runtime → 共享状态层：右侧总线通道（框外，不穿任何框）
bus_x = L + WM + 0.4
arrow(L + WM, 2.56, bus_x, 2.56, plain=True, color=BLUE, width=1.5)     # Host 出线
arrow(L + WM, 3.89, bus_x, 3.89, plain=True, color=BLUE, width=1.5)     # Runtime 出线
arrow(bus_x, 2.56, bus_x, 5.74, both=True, color=BLUE, width=1.5)       # 总线 → 状态层
bus_label = text(bus_x - 0.36, 4.1, 0.72, 0.6, ['PG 队列', '会话日志', 'NOTIFY'], size=7.5, color=BLUE, align=PP_ALIGN.CENTER, line_spacing=1.1)
bus_label.fill.solid(); bus_label.fill.fore_color.rgb = rgb('FFFFFF')   # 白底盖住穿过的总线段，标签"骑"在线上

# ── 架构特点（右）────────────────────────────────────────────────────
R, RW = 8.85, 4.0
panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(R), Inches(1.2), Inches(RW), Inches(6.16))
panel.adjustments[0] = 0.04
panel.fill.solid(); panel.fill.fore_color.rgb = rgb('FFFFFF')
panel.line.color.rgb = rgb(LINE); panel.line.width = Pt(1)
panel.shadow.inherit = False
text(R + 0.15, 1.28, RW - 0.3, 0.35, [[('架构特点', {'bold': True, 'size': 14, 'color': INK})]])
features = [
    ('全角色水平扩展', 'Host / Runtime / Collector 都是无状态 Deployment，KEDA 按连接数、队列深度自动扩缩（Runtime 池 2 → 20+）。数据面按节点分片，可管理成百上千个数据库节点、支撑数千人同时在线，扩容只是加副本。'),
    ('上下文与记忆跨 Pod 无缝切换', '会话日志、排队、任务、长期记忆全部落 PostgreSQL；任意 Runtime 都能接续任意会话，Host 多副本经 PG NOTIFY 扇出，浏览器落到哪个副本看到的都一致。Pod 滚动或故障时轮次自动重投、心跳所有权栅栏保证不丢提问、不重复执行。'),
    ('确定性采集 + 模型解读分离', '采集器产出确定性数字（榜单、占比、阈值判定、耗时构成、等待事件）直读入任务大盘，模型只做归因与建议；平台只读定位，能做什么由数据库账号权限决定。'),
    ('万物皆插件', '任务类型（健康检查 / Top SQL / WDR / DDL）、诊断工具、任务面板、存储后端均为插件；新增能力只加插件包，不改平台核心，双侧（服务端 + 浏览器端）各自注册。'),
]
y = 1.72
for i, (t, body) in enumerate(features):
    num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(R + 0.18), Inches(y + 0.02), Inches(0.3), Inches(0.3))
    num.fill.solid(); num.fill.fore_color.rgb = rgb(BLUE); num.line.fill.background(); num.shadow.inherit = False
    tf = num.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i + 1); r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = rgb('FFFFFF'); r.font.name = FONT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(R + 0.58, y - 0.02, RW - 0.75, 0.34, [[(t, {'bold': True, 'size': 12, 'color': INK})]])
    text(R + 0.58, y + 0.3, RW - 0.75, 1.1, [body], size=9.6, color=SUB, line_spacing=1.18)
    y += 1.4

text(0.5, 7.3, 8.0, 0.18, ['图例：实线箭头 = 直接连接；虚线 = 经共享状态层间接会合；×N = 可水平扩展副本'], size=7.5, color=DIM)
text(9.0, 7.3, 3.9, 0.18, ['opendb-harness · 架构一页 · 2026-08'], size=7.5, color=DIM, align=PP_ALIGN.RIGHT)

prs.save(OUT)
print('saved', OUT)
